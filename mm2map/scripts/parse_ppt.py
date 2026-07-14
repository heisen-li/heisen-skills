import sys
import os
import json

sys.stdout.reconfigure(encoding='utf-8')

def parse_pptx(pptx_path):
    try:
        from pptx import Presentation
    except ImportError:
        print(json.dumps({
            "error": "python-pptx not installed",
            "message": "Run: pip install python-pptx"
        }, ensure_ascii=False))
        sys.exit(1)

    if not os.path.exists(pptx_path):
        print(json.dumps({
            "error": "File not found",
            "message": f"PPT file not found: {pptx_path}"
        }, ensure_ascii=False))
        sys.exit(1)

    if not pptx_path.lower().endswith('.pptx'):
        print(json.dumps({
            "error": "Unsupported format",
            "message": "Only .pptx files are supported. .ppt (old format) needs conversion first."
        }, ensure_ascii=False))
        sys.exit(1)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(json.dumps({
            "error": "Failed to open PPT",
            "message": str(e)
        }, ensure_ascii=False))
        sys.exit(1)

    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": i,
            "title": "",
            "content": [],
            "level_structure": []
        }

        if slide.shapes.title_shape:
            slide_info["title"] = slide.shapes.title_shape.text.strip()

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level
                    slide_info["content"].append({
                        "text": text,
                        "level": level
                    })

                    if len(slide_info["level_structure"]) <= level:
                        while len(slide_info["level_structure"]) <= level:
                            slide_info["level_structure"].append([])
                    slide_info["level_structure"][level].append(text)

        if not slide_info["title"] and slide_info["content"]:
            first_text = slide_info["content"][0]["text"]
            slide_info["title"] = first_text[:80] if len(first_text) > 80 else first_text

        slides_data.append(slide_info)

    result = {
        "source": pptx_path,
        "total_slides": len(slides_data),
        "slides": slides_data
    }

    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({
            "error": "No PPT file path provided",
            "message": "Usage: python parse_ppt.py <pptx_file_path>"
        }, ensure_ascii=False))
        sys.exit(1)

    parse_pptx(sys.argv[1])
